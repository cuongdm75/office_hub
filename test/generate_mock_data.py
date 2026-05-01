import os
import sys
import subprocess

def install_deps():
    deps = ["openpyxl", "python-docx", "python-pptx"]
    subprocess.check_call([sys.executable, "-m", "pip", "install", *deps])

try:
    import openpyxl
    import docx
    import pptx
except ImportError:
    print("Installing dependencies...")
    install_deps()
    import openpyxl
    import docx
    import pptx

def generate_excel(path):
    print(f"Generating {path}...")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "DoanhThu_Q4"
    
    # Headers
    headers = ["Tháng", "Doanh Thu (VNĐ)", "Chi Phí (VNĐ)", "Lợi Nhuận (VNĐ)", "Đạt Chỉ Tiêu"]
    ws.append(headers)
    
    # Data
    data = [
        ["Tháng 10", 1500000000, 800000000, 700000000, "Yes"],
        ["Tháng 11", 1800000000, 900000000, 900000000, "Yes"],
        ["Tháng 12", 2500000000, 1200000000, 1300000000, "Yes"]
    ]
    for row in data:
        ws.append(row)
        
    wb.save(path)

def generate_word(path):
    print(f"Generating {path}...")
    doc = docx.Document()
    doc.add_heading("BÁO CÁO TỔNG KẾT DỰ ÁN NĂM 2025", 0)
    
    doc.add_heading("1. Tổng quan", level=1)
    doc.add_paragraph("Dự án Office Hub v1.1 đã hoàn thành xuất sắc các mục tiêu đề ra, bao gồm việc ra mắt Folder Scanner, cải thiện Outlook Agent và ứng dụng Mobile Remote UI.")
    
    doc.add_heading("2. Kết quả đạt được", level=1)
    doc.add_paragraph("Doanh thu quý 4 tăng trưởng 30% so với cùng kỳ năm trước. Hệ thống đã hoạt động ổn định không xảy ra lỗi critical nào.")
    
    doc.save(path)

def generate_ppt(path):
    print(f"Generating {path}...")
    prs = pptx.Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "OFFICE HUB V1.1 RELEASE"
    subtitle.text = "Sẵn sàng cho việc ra mắt"
    
    # Content slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide2.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Các tính năng mới"
    tf = body_shape.text_frame
    tf.text = "Hỗ trợ UIA, COM Automation"
    p = tf.add_paragraph()
    p.text = "Tích hợp Mobile React Native"
    p.level = 1
    
    prs.save(path)

if __name__ == "__main__":
    out_dir = r"e:\Office hub\test_mock_data"
    os.makedirs(out_dir, exist_ok=True)
    
    generate_excel(os.path.join(out_dir, "BaoCao_Q4.xlsx"))
    generate_word(os.path.join(out_dir, "DuAn_TongHop.docx"))
    generate_ppt(os.path.join(out_dir, "Presentation_Template.pptx"))
    
    print("Mock data generated successfully.")
